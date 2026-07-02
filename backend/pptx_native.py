"""Convert Figma frames into *editable* PowerPoint slides.

Strategy (hybrid native renderer):
  - TEXT             -> native PowerPoint text box, with per-character runs so a
                        bold/colored span inside one layer stays styled
  - RECTANGLE/ELLIPSE with a solid/no fill -> native autoshape (editable)
  - LINE             -> native line
  - Container fills   -> native background / slide background
  - Everything else (vectors, icons, gradients, image fills) -> rasterized PNG
    of just that node, pinned in place so the layout still looks right.

Coordinates are flattened to absolute slide positions (PowerPoint has no real
equivalent of Figma nesting / auto-layout). Content rotated, semi-transparent,
or off-canvas is handled best-effort."""

import io
import json
import math
import hashlib
import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

import figma

PX_TO_EMU = 9525  # 1 px @ 96 dpi
PX_TO_PT = 0.75   # 1 px -> 0.75 pt @ 96 dpi
SLIDE_W = Emu(12192000)  # 16:9
SLIDE_H = Emu(6858000)
MARGIN = Emu(228600)     # 0.25"

_ALIGN = {"LEFT": PP_ALIGN.LEFT, "CENTER": PP_ALIGN.CENTER,
          "RIGHT": PP_ALIGN.RIGHT, "JUSTIFIED": PP_ALIGN.JUSTIFY}
_ANCHOR = {"TOP": MSO_ANCHOR.TOP, "CENTER": MSO_ANCHOR.MIDDLE, "BOTTOM": MSO_ANCHOR.BOTTOM}
_CONTAINERS = {"FRAME", "GROUP", "COMPONENT", "INSTANCE", "SECTION", "COMPONENT_SET"}


# ---------- color / fill helpers ----------

def _solid(fills):
    """First visible solid paint -> (RGBColor, opacity 0..1), else None."""
    for p in fills or []:
        if p.get("visible", True) and p.get("type") == "SOLID":
            c = p["color"]
            opacity = p.get("opacity", 1) * c.get("a", 1)
            return RGBColor(round(c["r"] * 255), round(c["g"] * 255), round(c["b"] * 255)), opacity
    return None


def _has_complex_fill(fills):
    for p in fills or []:
        if p.get("visible", True) and (p.get("type", "").startswith("GRADIENT") or p.get("type") == "IMAGE"):
            return True
    return False


def _needs_raster(node):
    if node.get("type") in ("VECTOR", "STAR", "REGULAR_POLYGON", "BOOLEAN_OPERATION"):
        return True
    return _has_complex_fill(node.get("fills"))


def _has_text(node):
    """True if this subtree contains any non-empty TEXT (worth keeping editable)."""
    if node.get("visible", True) is False:
        return False
    if node.get("type") == "TEXT" and (node.get("characters") or "").strip():
        return True
    return any(_has_text(c) for c in node.get("children", []))


def _has_complex(node):
    """True if this subtree contains vectors / image / gradient fills (can't be native)."""
    if node.get("visible", True) is False:
        return False
    if node.get("type") in ("VECTOR", "STAR", "REGULAR_POLYGON", "BOOLEAN_OPERATION"):
        return True
    if _has_complex_fill(node.get("fills")):
        return True
    return any(_has_complex(c) for c in node.get("children", []))


def _set_alpha(fill_format, opacity):
    """Inject <a:alpha> into a solid fill so transparency carries over."""
    if opacity is None or opacity >= 0.999:
        return
    try:
        solid = fill_format._xPr.find(qn("a:solidFill"))
        srgb = solid.find(qn("a:srgbClr")) if solid is not None else None
        if srgb is None:
            return
        alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(max(0, min(1, opacity)) * 100000))})
        srgb.append(alpha)
    except Exception:  # noqa: BLE001
        pass


def _rotation_deg(node):
    t = node.get("relativeTransform")
    if not t or len(t) < 2:
        return 0.0
    try:
        return math.degrees(math.atan2(t[1][0], t[0][0]))
    except Exception:  # noqa: BLE001
        return 0.0


# ---------- public API ----------

def build_deck(file_key, token, frames, out_path, progress_cb=None):
    """frames: list of full node subtrees (from figma.fetch_nodes). Returns slide count."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = len(frames)
    done = 0

    cache = _RasterDiskCache(os.path.join(os.path.dirname(out_path), ".raster-cache"))

    for frame in frames:
        if not frame.get("absoluteBoundingBox"):
            done += 1
            if progress_cb:
                progress_cb(done, total)
            continue
        slide = prs.slides.add_slide(blank)
        _render_frame(slide, file_key, token, frame, cache)
        done += 1
        if progress_cb:
            progress_cb(done, total)

    prs.save(out_path)
    return len(prs.slides._sldIdLst)


def _render_frame(slide, file_key, token, frame, cache):
    fb = frame["absoluteBoundingBox"]
    fx, fy, fw, fh = fb["x"], fb["y"], fb["width"], fb["height"]
    avail_w = int(SLIDE_W) - 2 * int(MARGIN)
    avail_h = int(SLIDE_H) - 2 * int(MARGIN)
    fit = min(avail_w / (fw * PX_TO_EMU), avail_h / (fh * PX_TO_EMU))
    emu_scale = PX_TO_EMU * fit
    pt_scale = PX_TO_PT * fit
    ox = int(MARGIN) + (avail_w - fw * emu_scale) / 2
    oy = int(MARGIN) + (avail_h - fh * emu_scale) / 2

    def xywh(x, y, w, h):
        return (Emu(int(ox + (x - fx) * emu_scale)),
                Emu(int(oy + (y - fy) * emu_scale)),
                Emu(max(1, int(w * emu_scale))),
                Emu(max(1, int(h * emu_scale))))

    # Frame's own solid fill -> slide background (removes the letterbox border).
    root_fill = _solid(frame.get("fills"))
    if root_fill and not _has_complex_fill(frame.get("fills")):
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = root_fill[0]

    ops = []
    _walk(frame, ops, fb, is_root=True)

    pngs = {}
    missing_ids = []
    missing_hashes = {}
    for kind, node in ops:
        if kind != "raster":
            continue
        nid = node["id"]
        key = _node_hash(node)
        cached = cache.get(key)
        if cached is not None:
            pngs[nid] = cached
        else:
            missing_ids.append(nid)
            missing_hashes[nid] = key

    if missing_ids:
        rendered = figma.render_frames(file_key, token, missing_ids)
        for nid, data in rendered.items():
            pngs[nid] = data
            cache.set(missing_hashes[nid], data)

    for kind, node in ops:
        try:
            if kind == "raster":
                png = pngs.get(node["id"])
                if png:
                    _add_clipped_picture(slide, png, node["absoluteBoundingBox"], fb, xywh)
            elif kind == "text":
                _add_text(slide, node, xywh, pt_scale)
            elif kind == "rect":
                _add_rect(slide, node, xywh)
            elif kind == "ellipse":
                _add_ellipse(slide, node, xywh)
            elif kind == "line":
                _add_line(slide, node, xywh)
        except Exception:  # noqa: BLE001 — one bad node must never break the deck
            continue


def _xywh_of(node):
    b = node["absoluteBoundingBox"]
    return b["x"], b["y"], b["width"], b["height"]


def _node_hash(node: dict) -> str:
    blob = json.dumps(node, sort_keys=True, separators=(",", ":"), ensure_ascii=True)
    return hashlib.sha256(blob.encode("utf-8")).hexdigest()


class _RasterDiskCache:
    def __init__(self, root: str):
        self.root = root
        os.makedirs(self.root, exist_ok=True)

    def _path(self, key: str) -> str:
        return os.path.join(self.root, f"{key}.png")

    def get(self, key: str):
        p = self._path(key)
        if not os.path.exists(p):
            return None
        try:
            with open(p, "rb") as f:
                return f.read()
        except OSError:
            return None

    def set(self, key: str, data: bytes):
        p = self._path(key)
        try:
            with open(p, "wb") as f:
                f.write(data)
        except OSError:
            return


def _outside(b, root, tol=2):
    return (b["x"] > root["x"] + root["width"] + tol or b["x"] + b["width"] < root["x"] - tol or
            b["y"] > root["y"] + root["height"] + tol or b["y"] + b["height"] < root["y"] - tol)


def _walk(node, ops, root_bb, is_root=False):
    """Depth-first, back-to-front. Append (kind, node) draw ops in z-order."""
    if not is_root and node.get("visible", True) is False:
        return
    bb = node.get("absoluteBoundingBox")
    if not is_root and bb and _outside(bb, root_bb):
        return  # off-canvas / clipped content

    t = node.get("type")

    # If the whole frame is image/gradient-backed, render it flat (one image).
    if is_root and bb and _has_complex_fill(node.get("fills")):
        ops.append(("raster", node))
        return

    # Collapse a text-free graphic (icon, illustration, photo, icon cluster) into
    # ONE image instead of exploding it into dozens of vector pieces. Text-bearing
    # subtrees are recursed so their text/shapes stay editable.
    if not is_root and bb and not _has_text(node) and _has_complex(node):
        ops.append(("raster", node))
        return

    # A nested frame that CLIPS its content and contains graphics is a UI screen /
    # mockup (e.g. a phone screen). We can't honor Figma's clipping natively, so
    # extracting its inner text would leak hidden/clipped layers (nav drawers,
    # off-screen lists) onto the slide. Render the whole thing as one clean image.
    if not is_root and bb and node.get("clipsContent") and _has_complex(node):
        ops.append(("raster", node))
        return

    if t in _CONTAINERS or is_root:
        fills = node.get("fills")
        if bb and not is_root and _solid(fills) and not _has_complex_fill(fills):
            ops.append(("rect", node))  # native card / panel background
        for child in node.get("children", []):
            _walk(child, ops, root_bb)
        return

    if not bb:
        return
    if _needs_raster(node):
        ops.append(("raster", node))
    elif t == "TEXT":
        ops.append(("text", node))
    elif t == "RECTANGLE":
        ops.append(("rect", node))
    elif t == "ELLIPSE":
        ops.append(("ellipse", node))
    elif t == "LINE":
        ops.append(("line", node))
    else:
        ops.append(("raster", node))


# ---------- geometry with rotation ----------

def _box(node, xywh):
    """Return (left, top, width, height, rotation_deg). Uses the unrotated size +
    center when the node is rotated, so PowerPoint rotates it about its center."""
    rot = _rotation_deg(node)
    size = node.get("size")
    b = node["absoluteBoundingBox"]
    if abs(rot) > 0.5 and size and size.get("x") and size.get("y"):
        cx, cy = b["x"] + b["width"] / 2, b["y"] + b["height"] / 2
        left, top, width, height = xywh(cx - size["x"] / 2, cy - size["y"] / 2, size["x"], size["y"])
        return left, top, width, height, rot
    left, top, width, height = xywh(b["x"], b["y"], b["width"], b["height"])
    return left, top, width, height, 0.0


# ---------- native shape builders ----------

def _apply_fill_line(shape, node):
    fill = _solid(node.get("fills"))
    node_opacity = node.get("opacity", 1)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill[0]
        _set_alpha(shape.fill, fill[1] * node_opacity)
    else:
        shape.fill.background()

    stroke = _solid(node.get("strokes"))
    if stroke and node.get("strokeWeight"):
        shape.line.color.rgb = stroke[0]
        shape.line.width = Pt(node["strokeWeight"] * PX_TO_PT)
    else:
        shape.line.fill.background()

    try:
        shape.shadow.inherit = False
    except Exception:  # noqa: BLE001
        pass


def _add_rect(slide, node, xywh):
    left, top, width, height, rot = _box(node, xywh)
    radius = node.get("cornerRadius") or 0
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, left, top, width, height)
    if radius:
        short = min(node["absoluteBoundingBox"]["width"], node["absoluteBoundingBox"]["height"])
        try:
            shape.adjustments[0] = max(0.0, min(0.5, radius / short)) if short else 0.0
        except Exception:  # noqa: BLE001
            pass
    if rot:
        shape.rotation = rot
    _apply_fill_line(shape, node)


def _add_ellipse(slide, node, xywh):
    left, top, width, height, rot = _box(node, xywh)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    if rot:
        shape.rotation = rot
    _apply_fill_line(shape, node)


def _add_line(slide, node, xywh):
    left, top, width, height, _ = _box(node, xywh)
    conn = slide.shapes.add_connector(2, left, top, left + width, top + height)
    stroke = _solid(node.get("strokes"))
    if stroke:
        conn.line.color.rgb = stroke[0]
    conn.line.width = Pt((node.get("strokeWeight") or 1) * PX_TO_PT)


def _runs(node):
    """Split a TEXT node into paragraphs of styled runs, honoring per-character
    overrides. Returns [[(text, style_dict), ...], ...]."""
    chars = node.get("characters", "") or ""
    base = node.get("style", {})
    base_color = (_solid(node.get("fills")) or (RGBColor(0, 0, 0), 1))[0]
    overrides = node.get("characterStyleOverrides") or []
    table = node.get("styleOverrideTable") or {}

    def eff(i):
        key = overrides[i] if i < len(overrides) else 0
        ov = table.get(str(key), {}) if key else {}
        col = base_color
        if ov.get("fills"):
            s = _solid(ov["fills"])
            if s:
                col = s[0]
        return (
            ov.get("fontFamily", base.get("fontFamily", "Arial")),
            ov.get("fontSize", base.get("fontSize", 16)),
            (ov.get("fontWeight", base.get("fontWeight", 400)) or 400) >= 600,
            bool(ov.get("italic", base.get("italic", False))),
            col,
        )

    paragraphs, cur, buf, buf_style = [], [], "", None

    def flush():
        nonlocal buf
        if buf:
            cur.append((buf, buf_style))
            buf = ""

    for i, ch in enumerate(chars):
        if ch == "\n":
            flush()
            paragraphs.append(cur)
            cur, buf_style = [], None
            continue
        st = eff(i)
        if buf_style is None:
            buf_style = st
        elif st != buf_style:
            flush()
            buf_style = st
        buf += ch
    flush()
    paragraphs.append(cur)
    return paragraphs


def _add_text(slide, node, xywh, pt_scale):
    left, top, width, height, rot = _box(node, xywh)
    tb = slide.shapes.add_textbox(left, top, width, height)
    if rot:
        tb.rotation = rot
    tf = tb.text_frame
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, attr, Emu(0))

    style = node.get("style", {})
    chars = node.get("characters", "") or ""
    bb = node["absoluteBoundingBox"]
    lh_px = style.get("lineHeightPx") or (style.get("fontSize", 16) * 1.2)
    figma_lines = max(1, round(bb["height"] / lh_px)) if lh_px else 1

    # Single-line headings/labels must NOT wrap. PowerPoint's substituted fonts
    # are often wider than Figma's, so a one-line heading would wrap to two lines
    # and overflow onto the element below. Let it extend sideways instead.
    tf.word_wrap = not (figma_lines <= 1 and "\n" not in chars)

    tf.vertical_anchor = _ANCHOR.get(style.get("textAlignVertical", "TOP"), MSO_ANCHOR.TOP)
    align = _ALIGN.get(style.get("textAlignHorizontal", "LEFT"), PP_ALIGN.LEFT)

    # Use Figma's EXACT line height (in points), not a multiple. A multiple gets
    # compounded with the font's own leading in PowerPoint, making every block
    # ~20% taller and overflowing onto the element below it.
    line_pt = Pt(round(style["lineHeightPx"] * pt_scale, 2)) if style.get("lineHeightPx") else None

    for i, para in enumerate(_runs(node)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        if line_pt is not None:
            p.line_spacing = line_pt
        if not para:  # blank line
            continue
        for text, (family, size_px, bold, italic, color) in para:
            run = p.add_run()
            run.text = text
            f = run.font
            f.size = Pt(max(1, round(size_px * pt_scale)))
            f.name = family
            f.bold = bold
            f.italic = italic
            f.color.rgb = color


def _add_clipped_picture(slide, png, nb, fb, xywh):
    """Place a rendered node image, clipped to the frame bounds — Figma renders a
    node unclipped, but the frame has clipsContent, so overflow must be cut off
    instead of spilling past the slide edge."""
    nx, ny, nw, nh = nb["x"], nb["y"], nb["width"], nb["height"]
    fx, fy, fw, fh = fb["x"], fb["y"], fb["width"], fb["height"]
    ix0, iy0 = max(nx, fx), max(ny, fy)
    ix1, iy1 = min(nx + nw, fx + fw), min(ny + nh, fy + fh)
    if ix1 <= ix0 or iy1 <= iy0 or nw <= 0 or nh <= 0:
        return  # fully outside the frame

    left, top, width, height = xywh(ix0, iy0, ix1 - ix0, iy1 - iy0)
    pic = slide.shapes.add_picture(io.BytesIO(png), left, top, width, height)

    # Crop away the portion of the source image that falls outside the frame.
    for attr, val in (
        ("crop_left", (ix0 - nx) / nw),
        ("crop_right", (nx + nw - ix1) / nw),
        ("crop_top", (iy0 - ny) / nh),
        ("crop_bottom", (ny + nh - iy1) / nh),
    ):
        if val > 0.0001:
            setattr(pic, attr, round(min(val, 0.99), 5))
